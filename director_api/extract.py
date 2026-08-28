"""Extract text and structured brief fields from any uploaded format."""

from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path

try:
    import yaml
except ImportError:  # pragma: no cover
    yaml = None

KNOWN_FIELDS = (
    "brand",
    "product",
    "therapy_area",
    "indication",
    "market",
    "business_goal",
    "target_specialties",
    "hcp_segments",
    "brand_evidence",
    "existing_evidence",
    "evolving_evidence",
    "guidelines",
    "hcp_insights",
    "competitors",
    "access_and_cost",
    "constraints",
    "notes",
)

TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".csv",
    ".tsv",
    ".json",
    ".yaml",
    ".yml",
    ".html",
    ".htm",
    ".xml",
    ".rtf",
    ".log",
    ".outline",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
OFFICE_SUFFIXES = {".docx", ".pptx", ".xlsx", ".xls", ".odt", ".odp", ".ods", ".ppt", ".doc"}
PDF_SUFFIXES = {".pdf"}

SUPPORTED_SUFFIXES = TEXT_SUFFIXES | IMAGE_SUFFIXES | OFFICE_SUFFIXES | PDF_SUFFIXES | {".key"}


@dataclass
class ExtractedFile:
    filename: str
    mime: str
    suffix: str
    text: str
    notes: list[str] = field(default_factory=list)
    pages: int | None = None
    bytes: int = 0


@dataclass
class ExtractedBrief:
    brand: str = ""
    product: str = ""
    therapy_area: str = ""
    indication: str = ""
    market: str = ""
    business_goal: str = ""
    target_specialties: list[str] = field(default_factory=list)
    hcp_segments: list[str] = field(default_factory=list)
    brand_evidence: list[str] = field(default_factory=list)
    existing_evidence: list[str] = field(default_factory=list)
    evolving_evidence: list[str] = field(default_factory=list)
    guidelines: list[str] = field(default_factory=list)
    hcp_insights: list[str] = field(default_factory=list)
    competitors: list[str] = field(default_factory=list)
    access_and_cost: list[str] = field(default_factory=list)
    constraints: list[str] = field(default_factory=list)
    notes: str = ""
    source_files: list[str] = field(default_factory=list)
    raw_text: str = ""
    extraction_notes: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return asdict(self)


def extract_files(uploads: list[tuple[str, bytes, str]]) -> list[ExtractedFile]:
    """Parse many uploads. Each tuple is (filename, payload, mime)."""
    return [extract_one(name, payload, mime) for name, payload, mime in uploads]


def is_design_template(filename: str, text: str = "") -> bool:
    """True for a visual master (Omnicom PPTX), not a client brief."""
    name = (filename or "").lower()
    if "template" in name or "omnicom" in name:
        return True
    blob = (text or "").lower()
    if "omnicom group" in blob and "medical strategy deck" in blob:
        return True
    return False


def partition_uploads(
    uploads: list[tuple[str, bytes, str]],
) -> tuple[list[tuple[str, bytes, str]], list[tuple[str, bytes, str]]]:
    """Split design-template files from brief files before extraction."""
    templates: list[tuple[str, bytes, str]] = []
    briefs: list[tuple[str, bytes, str]] = []
    for item in uploads:
        name, payload, mime = item
        suffix = Path(name).suffix.lower()
        preview = ""
        if suffix == ".pptx":
            try:
                preview = extract_one(name, payload, mime).text
            except Exception:
                preview = ""
        if is_design_template(name, preview):
            templates.append(item)
        else:
            briefs.append(item)
    return templates, briefs


def extract_one(filename: str, payload: bytes, mime: str = "") -> ExtractedFile:
    suffix = Path(filename).suffix.lower()
    notes: list[str] = []
    text = ""
    pages = None

    try:
        if suffix in {".yaml", ".yml"}:
            text = _from_yaml(payload)
        elif suffix == ".json":
            text = _from_json(payload)
        elif suffix in {".csv", ".tsv"}:
            text = _from_csv(payload, "\t" if suffix == ".tsv" else ",")
        elif suffix in {".html", ".htm", ".xml"}:
            text = _strip_markup(_decode(payload))
        elif suffix == ".rtf":
            text = _from_rtf(_decode(payload))
        elif suffix == ".pdf":
            text, pages = _from_pdf(payload)
        elif suffix == ".docx":
            text = _from_docx(payload)
        elif suffix == ".pptx":
            text = _from_pptx(payload)
        elif suffix in {".xlsx", ".xls"}:
            text = _from_xlsx(payload)
        elif suffix in {".odt", ".odp", ".ods"}:
            text = _from_opendocument(payload)
        elif suffix in {".ppt", ".doc"}:
            text = _from_ole_legacy(payload)
            if not text.strip():
                notes.append(
                    f"{suffix} is a legacy Office binary. Convert to .{suffix[1:]}x for richer extraction."
                )
        elif suffix in IMAGE_SUFFIXES:
            notes.append(
                "Image uploaded. Text was not OCR'd locally; attach a text/PDF brief for full extraction, "
                "or keep the image as a visual reference in the working file."
            )
            text = f"[Image attachment: {filename}]"
        elif suffix in TEXT_SUFFIXES or not suffix:
            text = _decode(payload)
        else:
            text = _decode(payload)
            if not text.strip():
                notes.append(f"No text extractor for {suffix or 'unknown type'}; stored as raw attachment.")
    except Exception as exc:  # keep the pipeline moving; record the failure
        notes.append(f"Partial extract for {filename}: {exc}")
        text = _decode(payload)

    return ExtractedFile(
        filename=filename,
        mime=mime or "application/octet-stream",
        suffix=suffix,
        text=text.strip(),
        notes=notes,
        pages=pages,
        bytes=len(payload),
    )


def merge_into_brief(files: list[ExtractedFile], pasted: str = "") -> ExtractedBrief:
    """Turn extracted file text + optional pasted copy into a structured brief."""
    chunks = [f.text for f in files if f.text]
    if pasted.strip():
        chunks.append(pasted.strip())
    raw = "\n\n---\n\n".join(chunks)
    brief = _parse_structured(raw)
    brief.raw_text = raw
    brief.source_files = [f.filename for f in files]
    brief.extraction_notes = [n for f in files for n in f.notes]
    _infer_missing(brief)
    return brief


def _decode(payload: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return payload.decode(enc)
        except UnicodeDecodeError:
            continue
    return payload.decode("utf-8", errors="replace")


def _from_yaml(payload: bytes) -> str:
    text = _decode(payload)
    if yaml is None:
        return text
    data = yaml.safe_load(text)
    if isinstance(data, dict):
        return _mapping_to_markdown(data)
    return text


def _from_json(payload: bytes) -> str:
    data = json.loads(_decode(payload))
    if isinstance(data, dict):
        return _mapping_to_markdown(data)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _from_csv(payload: bytes, delimiter: str) -> str:
    buf = io.StringIO(_decode(payload))
    reader = csv.reader(buf, delimiter=delimiter)
    rows = list(reader)
    if not rows:
        return ""
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows if any(row))


def _from_pdf(payload: bytes) -> tuple[str, int | None]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(payload))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        content = page.extract_text() or ""
        if content.strip():
            pages.append(f"## Page {i}\n{content}")
    return "\n\n".join(pages), len(reader.pages)


def _from_docx(payload: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(payload))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_pptx(payload: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(payload))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        bits = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    bits.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        bits.append(" | ".join(cells))
        if bits:
            slides.append(f"## Slide {i}\n" + "\n".join(bits))
    return "\n\n".join(slides)


def _from_xlsx(payload: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(payload), data_only=True, read_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c not in (None, "")]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"## Sheet {name}\n" + "\n".join(rows[:400]))
    return "\n\n".join(sheets)


def _from_opendocument(payload: bytes) -> str:
    """Read text from ODF packages (content.xml)."""
    with zipfile.ZipFile(io.BytesIO(payload)) as zf:
        xml = zf.read("content.xml").decode("utf-8", errors="replace")
    return _strip_markup(xml)


def _from_ole_legacy(payload: bytes) -> str:
    """Best-effort string harvest from old .doc / .ppt binaries."""
    strings = re.findall(rb"[\x20-\x7e]{6,}", payload)
    text = "\n".join(s.decode("ascii", errors="ignore") for s in strings)
    return "\n".join(line for line in text.splitlines() if not _noise_line(line))


def _from_rtf(text: str) -> str:
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ")
    return re.sub(r"\s+", " ", text).strip()


def _strip_markup(text: str) -> str:
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def _noise_line(line: str) -> bool:
    s = line.strip()
    if len(s) < 4:
        return True
    return bool(re.fullmatch(r"[A-Z0-9_\-.]{4,}", s)) and " " not in s


def _mapping_to_markdown(data: dict) -> str:
    lines = ["# CLIENT BRIEF", ""]
    for key, value in data.items():
        lines.append(f"## {str(key).replace('_', ' ').title()}")
        if isinstance(value, list):
            lines.extend(f"- {item}" for item in value)
        elif isinstance(value, dict):
            for k, v in value.items():
                lines.append(f"- **{k}**: {v}")
        else:
            lines.append(str(value))
        lines.append("")
    return "\n".join(lines)


def _parse_structured(raw: str) -> ExtractedBrief:
    """Accept YAML/JSON-shaped text or headed Markdown and fill known fields."""
    brief = ExtractedBrief()
    parsed = _try_load_mapping(raw)
    if parsed:
        _apply_mapping(brief, parsed)
        _tidy_brief(brief, raw)
        return brief

    current = None
    bucket: list[str] = []

    def flush():
        nonlocal current, bucket
        if current and current in KNOWN_FIELDS:
            _assign(brief, current, bucket)
        current = None
        bucket = []

    for line in raw.splitlines():
        stripped = line.strip()
        heading = re.match(r"^#{1,3}\s+(.+)$", stripped)
        label = re.match(
            r"^([A-Za-z][A-Za-z0-9 /&()_,'-]{1,70})\s*:\s*(.*)$",
            stripped,
        )
        key = None
        rest = ""
        if heading:
            key = _normalize_key(heading.group(1))
        elif label:
            key = _normalize_key(label.group(1))
            rest = label.group(2).strip()
        else:
            prefix = _section_prefix(stripped)
            if prefix:
                key = prefix
        if key and key in KNOWN_FIELDS:
            flush()
            current = key
            if rest:
                bucket.append(rest)
            continue
        if label and key not in KNOWN_FIELDS:
            title = label.group(1)
            heading_like = (
                len(title) < 42
                and "," not in title
                and title[:1].isupper()
                and not title.lower().startswith(("consultant", "indian", "named"))
            )
            if heading_like:
                flush()
                current = None
                continue
        if current:
            cleaned = stripped.lstrip("-*•").strip()
            if cleaned:
                bucket.append(cleaned)
    flush()
    _tidy_brief(brief, raw)
    return brief


def _try_load_mapping(raw: str) -> dict | None:
    text = raw.strip()
    if text.startswith("{") and text.endswith("}"):
        try:
            data = json.loads(text)
            return data if isinstance(data, dict) else None
        except json.JSONDecodeError:
            return None
    if yaml is None or not re.search(r"^\w+:", text, re.M):
        return None
    try:
        data = yaml.safe_load(text)
    except Exception:
        return None
    if not isinstance(data, dict):
        return None
    known = sum(1 for k in data if _normalize_key(str(k)) in KNOWN_FIELDS)
    if known < 2:
        return None
    snake = sum(1 for k in data if re.fullmatch(r"[a-z][a-z0-9_]*", str(k) or ""))
    if snake < 2:
        # Title-case "Brand:" client briefs are not YAML. Parse them as labeled prose.
        return None
    brand = str(data.get("brand") or data.get("Brand") or "")
    if len(brand) > 80 or "molecule" in brand.lower() or "\n" in brand:
        return None
    return data


def _apply_mapping(brief: ExtractedBrief, data: dict) -> None:
    for key, value in data.items():
        norm = _normalize_key(str(key))
        if norm in KNOWN_FIELDS:
            _assign(brief, norm, value)


def _assign(brief: ExtractedBrief, key: str, value) -> None:
    list_fields = {
        "target_specialties",
        "hcp_segments",
        "brand_evidence",
        "existing_evidence",
        "evolving_evidence",
        "guidelines",
        "hcp_insights",
        "competitors",
        "access_and_cost",
        "constraints",
    }
    if key in list_fields:
        items = _as_list(value)
        setattr(brief, key, items)
    elif key == "notes":
        brief.notes = _as_text(value)
    else:
        setattr(brief, key, _as_text(value))


def _as_list(value) -> list[str]:
    if value is None:
        return []
    if isinstance(value, list):
        out = []
        for item in value:
            if isinstance(item, (dict, list)):
                out.append(json.dumps(item, ensure_ascii=False))
            else:
                text = str(item).strip()
                if text:
                    out.append(text)
        return _merge_wrapped(out)
    text = str(value).strip()
    if not text:
        return []
    if "\n" in text:
        return _merge_wrapped([ln.lstrip("-*• ").strip() for ln in text.splitlines() if ln.strip()])
    return [text]


def _as_text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, list):
        joined = " ".join(str(v).strip() for v in value if str(v).strip())
        return re.sub(r"\s+", " ", joined.replace(" ; ", " ")).strip()
    return str(value).strip()


def _merge_wrapped(items: list[str]) -> list[str]:
    """Rejoin a client brief that wrapped mid-sentence; drop heading leftovers."""
    out: list[str] = []
    for raw in items:
        s = re.sub(r"\s+", " ", (raw or "").strip())
        if not s:
            continue
        # Wrap of a heading that already opened the section: "consultant physicians…):"
        if re.search(r"\)\s*:\s*$", s) and not re.match(r"^\d+", s) and len(s) < 90:
            continue
        numbered = bool(re.match(r"^\d+[\.)]\s+", s))
        bulleted = bool(re.match(r"^[-*•]\s+", s))
        body = re.sub(r"^(\d+[\.)]\s+|[-*•]\s+)", "", s).strip() or s
        if out and not numbered and not bulleted and s[:1].islower():
            out[-1] = f"{out[-1]} {body}".strip()
        else:
            out.append(body)
    return out


def _normalize_key(label: str) -> str:
    s = re.sub(r"[^a-z0-9]+", "_", label.lower()).strip("_")
    s = re.sub(r"_+", "_", s).strip("_")
    aliases = {
        "therapy": "therapy_area",
        "ta": "therapy_area",
        "therapy_area": "therapy_area",
        "therapeutic_area": "therapy_area",
        "brand_name": "brand",
        "product_name": "product",
        "molecule": "product",
        "molecule_composition": "product",
        "composition": "product",
        "goal": "business_goal",
        "objective": "business_goal",
        "business_objective": "business_goal",
        "business_objective_what_must_move_by_when_by_how_much": "business_goal",
        "specialties": "target_specialties",
        "target_specialties": "target_specialties",
        "targets": "target_specialties",
        "segments": "hcp_segments",
        "hcp_tiers_status_mix": "hcp_segments",
        "evidence": "brand_evidence",
        "evidence_in_hand": "brand_evidence",
        "evidence_in_hand_attached_see_test_source_pack_md": "brand_evidence",
        "insights": "hcp_insights",
        "hcp_insight": "hcp_insights",
        "in_house_hcp_insights": "hcp_insights",
        "competitor": "competitors",
        "competitive_set": "competitors",
        "competitive_set_what_the_narrative_must_displace": "competitors",
        "access": "access_and_cost",
        "cost": "access_and_cost",
        "price_and_access": "access_and_cost",
        "price_and_access_context": "access_and_cost",
        "constraint": "constraints",
        "mlr_constraints": "constraints",
        "guidelines_in_play": "guidelines",
        "evolving_new_evidence": "evolving_evidence",
        "country": "market",
        "geography": "market",
        "geography_and_city_tiers": "market",
    }
    if s in aliases:
        return aliases[s]
    for prefix, dest in (
        ("business_objective", "business_goal"),
        ("in_house_hcp", "hcp_insights"),
        ("hcp_insight", "hcp_insights"),
        ("competitive_set", "competitors"),
        ("evidence_in_hand", "brand_evidence"),
        ("guidelines", "guidelines"),
        ("price_and_access", "access_and_cost"),
        ("mlr", "constraints"),
        ("molecule", "product"),
        ("therapeutic", "therapy_area"),
        ("target_specialt", "target_specialties"),
        ("evolving", "evolving_evidence"),
    ):
        if s.startswith(prefix):
            return dest
    return s


def _infer_missing(brief: ExtractedBrief) -> None:
    """Fill obvious gaps from free text when structured fields are empty."""
    blob = brief.raw_text
    if not brief.brand:
        m = re.search(r"\bbrand\s*[:\-]\s*([A-Z][\w\- ]{1,40})", blob, re.I)
        if m:
            brief.brand = m.group(1).strip()
    if not brief.therapy_area:
        m = re.search(
            r"(cardiology|oncology|diabetes|neurology|respiratory|immunology|"
            r"dermatology|infectious|heart failure|HFrEF|HFpEF|COPD)[^\n]{0,40}",
            blob,
            re.I,
        )
        if m:
            brief.therapy_area = m.group(0).strip()
    if not brief.market:
        m = re.search(r"\b(India|United States|USA|UK|EU|China|Brazil|Japan|Germany|France)\b", blob)
        if m:
            brief.market = m.group(1)
    _tidy_brief(brief, blob)


def _section_prefix(line: str) -> str | None:
    low = line.lower()
    pairs = (
        ("in-house hcp insight", "hcp_insights"),
        ("hcp insight", "hcp_insights"),
        ("guidelines in play", "guidelines"),
        ("evidence in hand", "brand_evidence"),
        ("competitive set", "competitors"),
        ("business objective", "business_goal"),
        ("price and access", "access_and_cost"),
        ("mlr constraint", "constraints"),
        ("target specialt", "target_specialties"),
        ("molecule / composition", "product"),
        ("therapeutic area", "therapy_area"),
    )
    for prefix, field in pairs:
        if low.startswith(prefix):
            return field
    return None


def _tidy_brief(brief: ExtractedBrief, raw: str = "") -> None:
    """Stop a labeled client brief from collapsing into one giant brand string."""
    raw = raw or brief.raw_text or ""
    brief.brand = _short_name(brief.brand, 48)
    if brief.product:
        brief.product = _clip_text(brief.product, 160)
    elif raw:
        mol = re.search(
            r"molecule[^\n:]{0,20}:\s*(.+?)(?:\n[A-Z][^\n:]{0,40}:|\n\n|$)",
            raw,
            re.I | re.S,
        )
        if mol:
            brief.product = _clip_text(re.sub(r"\s+", " ", mol.group(1)), 160)
    brief.therapy_area = _clip_text(brief.therapy_area, 80)
    brief.indication = _first_block(brief.indication, 280)
    brief.business_goal = _first_block(brief.business_goal, 280)
    if "molecule" in (brief.brand or "").lower() or len(brief.brand or "") > 48:
        m = re.search(r"\bbrand\s*:\s*([^\n]+)", raw, re.I)
        if m:
            brief.brand = _short_name(m.group(1), 48)


def _short_name(text: str, limit: int = 48) -> str:
    text = (text or "").replace("\r", " ").strip().split("\n")[0]
    text = re.split(r"\s*;\s*|molecule\s*/|therapeutic\s+area", text, maxsplit=1, flags=re.I)[0]
    text = re.sub(r"\s+", " ", text).strip(" ,;—-")
    return text[:limit].strip()


def _clip_text(text: str, limit: int) -> str:
    text = re.sub(r"\s+", " ", (text or "").replace("\r", " ")).strip()
    return text[:limit].strip()


def _first_block(text: str, limit: int) -> str:
    text = (text or "").replace("\r", "\n").strip()
    for stop in (
        "STATED POPULATION",
        "Business objective",
        "Competitive set",
        "Evidence in hand",
        "In-house HCP",
    ):
        idx = text.lower().find(stop.lower())
        if idx > 20:
            text = text[:idx].strip()
    para = text.split("\n\n")[0].strip()
    return _clip_text(para.replace("\n", " "), limit)
