"""Build an editable 16:9 PowerPoint from a STRATA strategy pack.

Shapes, text, tables, and native Office charts are written as real PPT
objects so a client can open the file in PowerPoint or Google Slides and
keep editing. Forest and box plots are drawn with shapes (still editable).
"""

from __future__ import annotations

import io
import re
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x11, 0x11, 0x11)
NAVY = RGBColor(0x00, 0x00, 0x00)
CREAM = RGBColor(0xF7, 0xF7, 0xF4)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
COPPER = RGBColor(0xFF, 0x64, 0x33)
TEAL = RGBColor(0x4E, 0x7D, 0xF2)
CRIMSON = RGBColor(0xFF, 0x64, 0x33)
MUTED = RGBColor(0x77, 0x77, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x4E, 0x7D, 0xF2)
ORANGE = RGBColor(0xFF, 0x64, 0x33)
BODY = RGBColor(0x33, 0x33, 0x33)
PAGE_GREY = RGBColor(0xA7, 0xA7, 0xA7)
DARK_CARD = RGBColor(0x11, 0x11, 0x11)
HAIRLINE = RGBColor(0x00, 0x00, 0x00)
FONT = "Calibri"
DISPLAY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def pack_to_pptx(pack: dict) -> bytes:
    """Return a .pptx byte string for the given strategy pack."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    meta = pack.get("meta") or {}
    doctrine = pack.get("doctrine") or {}
    slides = pack.get("slides") or []

    for spec in slides:
        slide = prs.slides.add_slide(blank)
        dark = spec.get("layout") in {"title", "close", "idea"}
        _fill_slide(slide, NAVY if dark else CREAM)
        _render_slide(slide, spec, dark)
        _notes(slide, spec)

    _appendix_bibliography(prs, blank, pack)
    _appendix_interventions(prs, blank, pack)
    _appendix_dashboard(prs, blank, pack)
    _appendix_edit_guide(prs, blank, meta, doctrine)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def filename_for(pack: dict) -> str:
    brand = (pack.get("meta") or {}).get("brand") or "strategy"
    slug = re.sub(r"[^A-Za-z0-9]+", "-", brand).strip("-") or "strategy"
    return f"{slug}-strategy-deck.pptx"


def _fill_slide(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    # Send the background behind other shapes.
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def _textbox(slide, l, t, w, h, text, *, size=16, bold=False, color=INK, font=FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box


def _as_text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, dict):
        return "; ".join(f"{k}: {v}" for k, v in value.items())
    if isinstance(value, list):
        return "; ".join(_as_text(v) for v in value)
    return str(value)


def _bullets(slide, l, t, w, h, items: list, *, color=INK, size=15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = _as_text(item)
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(8)
    return box


def _table(slide, l, t, w, h, headers: list[str], rows: list[list[str]]):
    cols = max(len(headers), 1)
    n_rows = 1 + len(rows)
    table_shape = slide.shapes.add_table(n_rows, cols, l, t, w, h)
    table = table_shape.table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        _style_cell(cell, NAVY, WHITE, bold=True, size=11)
    for i, row in enumerate(rows, 1):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = _as_text(row[j] if j < len(row) else "")
            bg = PAPER if i % 2 else CREAM
            _style_cell(cell, bg, INK, size=11)
    return table_shape


def _style_cell(cell, bg: RGBColor, fg: RGBColor, *, bold=False, size=12):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = fg
        p.font.name = FONT
        p.alignment = PP_ALIGN.LEFT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _notes(slide, spec: dict) -> None:
    bits = [spec.get("narrative") or ""]
    if spec.get("callout"):
        bits.append(f"{spec['callout'].get('label')}: {spec['callout'].get('text')}")
    chart = spec.get("chart") or {}
    if chart.get("note"):
        bits.append(chart["note"])
    notes = slide.notes_slide.notes_text_frame
    notes.text = "\n\n".join(b for b in bits if b)


def _hairline(slide, l, t, w, color=HAIRLINE) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Emu(6350))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _rail_box(slide, l, t, w, h, accent: RGBColor, fill: RGBColor = PAPER) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.06), h)
    rail.fill.solid()
    rail.fill.fore_color.rgb = accent
    rail.line.fill.background()


def _accent(i, card=None) -> RGBColor:
    if isinstance(card, dict) and card.get("accent") == "orange":
        return ORANGE
    if isinstance(card, dict) and card.get("accent") == "blue":
        return BLUE
    return BLUE if i % 2 == 0 else ORANGE


def _chrome(slide, spec: dict, dark: bool) -> None:
    _textbox(slide, Inches(12.82), Inches(0.18), Inches(0.4), Inches(0.18),
             spec.get("page") or "", size=8, color=PAGE_GREY, align=PP_ALIGN.RIGHT)
    source = spec.get("source") or spec.get("footnote") or ""
    if source:
        _textbox(slide, Inches(0.34), Inches(6.82), Inches(10.4), Inches(0.22),
                 source, size=9, color=RGBColor(0x80, 0x80, 0x80) if dark else MUTED)
    elif spec.get("refs"):
        from .cite import format_marks
        nums = [int(n) for n in spec["refs"] if n not in (None, "")]
        if nums:
            _textbox(slide, Inches(0.34), Inches(6.82), Inches(10.4), Inches(0.22),
                     f"References {format_marks(nums)}  ·  Vancouver list at end of deck",
                     size=9, color=MUTED)


def _stat_cards(slide, stats: list[dict], l, t, w, h) -> None:
    if not stats:
        return
    n = max(len(stats), 1)
    cw = w / n
    for i, st in enumerate(stats):
        x = l + cw * i + Inches(0.08)
        card_w = cw - Inches(0.16)
        _rail_box(slide, x, t, card_w, h, _accent(i, st), PAPER)
        _textbox(slide, x + Inches(0.18), t + Inches(0.16), card_w - Inches(0.28), Inches(0.55),
                 str(st.get("value") or ""), size=28, bold=True, color=INK, font=DISPLAY)
        _textbox(slide, x + Inches(0.18), t + Inches(0.78), card_w - Inches(0.28), h - Inches(0.92),
                 st.get("caption") or "", size=12, color=BODY)


def _render_slide(slide, spec: dict, dark: bool) -> None:
    ink = WHITE if dark else INK
    body = RGBColor(0xBF, 0xBF, 0xBF) if dark else BODY
    layout = spec.get("layout")
    _chrome(slide, spec, dark)

    if layout == "title":
        _textbox(slide, Inches(0.55), Inches(1.72), Inches(8.4), Inches(0.28),
                 (spec.get("kicker") or "").upper(), size=12, bold=True, color=BLUE)
        _textbox(slide, Inches(0.55), Inches(2.05), Inches(8.6), Inches(1.15),
                 spec.get("title") or "", size=32, bold=True, color=WHITE, font=DISPLAY)
        if spec.get("subtitle"):
            _textbox(slide, Inches(0.55), Inches(3.35), Inches(8.8), Inches(0.7),
                     spec["subtitle"], size=16, color=RGBColor(0xB7, 0xC5, 0xE8))
        if spec.get("narrative"):
            _textbox(slide, Inches(0.55), Inches(4.1), Inches(8.8), Inches(0.55),
                     spec["narrative"], size=13, color=RGBColor(0xBF, 0xBF, 0xBF))
        if spec.get("cards"):
            _dark_chips(slide, spec["cards"], Inches(0.55), Inches(4.85), Inches(12.2), Inches(1.55))
        elif spec.get("chart"):
            _add_chart(slide, spec["chart"], Inches(0.55), Inches(4.85), Inches(12.2), Inches(1.7))
        return

    if layout in {"close", "idea"}:
        _textbox(slide, Inches(0.55), Inches(1.2), Inches(8.0), Inches(0.22),
                 (spec.get("kicker") or "").upper(), size=10, bold=True, color=ORANGE)
        _textbox(slide, Inches(0.55), Inches(1.48), Inches(9.4), Inches(0.9),
                 spec.get("title") or "", size=22, bold=True, color=WHITE, font=DISPLAY)
        if spec.get("subtitle") or spec.get("narrative"):
            _textbox(slide, Inches(0.55), Inches(2.4), Inches(9.0), Inches(0.55),
                     spec.get("subtitle") or spec.get("narrative") or "", size=13, color=RGBColor(0xBF, 0xBF, 0xBF))
        if spec.get("cards"):
            _dark_quad(slide, spec["cards"], Inches(0.55), Inches(3.15), Inches(12.1), Inches(2.9))
        elif spec.get("chart"):
            _add_chart(slide, spec["chart"], Inches(0.55), Inches(3.2), Inches(12.2), Inches(3.2))
        if spec.get("callout"):
            _textbox(slide, Inches(0.55), Inches(6.2), Inches(10.5), Inches(0.35),
                     spec["callout"].get("text") or "", size=12, bold=True, color=ORANGE)
        return

    if not dark:
        _hairline(slide, Inches(0.34), Inches(0.55), Inches(11.95), HAIRLINE)
        _hairline(slide, Inches(0.34), Inches(1.42), Inches(1.36), HAIRLINE)

    _textbox(slide, Inches(0.34), Inches(0.82), Inches(8.4), Inches(0.22),
             (spec.get("kicker") or "").upper(), size=10, bold=True, color=ORANGE)
    _textbox(slide, Inches(0.34), Inches(1.04), Inches(8.6), Inches(0.7),
             spec.get("title") or "", size=18, bold=True, color=INK, font=DISPLAY)

    if layout == "insight":
        _textbox(slide, Inches(0.34), Inches(2.08), Inches(3.35), Inches(1.35),
                 spec.get("narrative") or spec.get("subtitle") or "", size=12, color=BODY)
        stats = spec.get("stats") or []
        if stats:
            _stat_cards(slide, stats, Inches(4.15), Inches(1.28), Inches(8.5), Inches(2.15))
        if spec.get("soWhat"):
            _textbox(slide, Inches(4.15), Inches(3.6), Inches(8.4), Inches(0.7),
                     spec["soWhat"], size=16, color=INK, font=DISPLAY)
        if spec.get("chart"):
            _add_chart(slide, spec["chart"], Inches(4.15), Inches(4.35), Inches(8.5), Inches(2.2))
        elif spec.get("table") and not stats:
            table = spec["table"]
            _table(slide, Inches(4.15), Inches(3.55), Inches(8.5), Inches(2.9),
                   table.get("headers") or [], table.get("rows") or [])
        return

    top = Inches(1.85)
    if spec.get("subtitle"):
        _textbox(slide, Inches(0.34), top, Inches(12.4), Inches(0.4),
                 spec["subtitle"], size=13, color=body)
        top += Inches(0.38)
    if spec.get("narrative"):
        _textbox(slide, Inches(0.34), top, Inches(12.4), Inches(0.45),
                 spec.get("narrative") or "", size=13, color=BODY)
        top += Inches(0.48)

    if layout == "split":
        if spec.get("bullets"):
            _bullets(slide, Inches(0.34), top, Inches(5.5), Inches(3.4),
                     spec["bullets"], color=INK, size=14)
        if spec.get("table"):
            table = spec["table"]
            _table(slide, Inches(0.34), top, Inches(5.7), Inches(3.6),
                   table.get("headers") or [], table.get("rows") or [])
        if spec.get("callout"):
            _callout(slide, Inches(0.34), Inches(6.25), Inches(5.7), Inches(0.5), spec["callout"], dark)
        if spec.get("chart"):
            _add_chart(slide, spec["chart"], Inches(6.3), top, Inches(6.5), Inches(4.5))
        elif spec.get("cards"):
            _cards(slide, spec["cards"], Inches(6.3), top, Inches(6.5), Inches(4.5))
        return

    if layout in {"infographic", "chart"}:
        if spec.get("chart"):
            _add_chart(slide, spec["chart"], Inches(0.34), top, Inches(12.6), Inches(6.55) - top)
        return

    if layout in {"references", "table"}:
        table = spec.get("table") or {}
        if spec.get("chart") and table.get("rows"):
            _table(slide, Inches(0.34), top, Inches(7.5), Inches(6.55) - top,
                   table.get("headers") or [], table.get("rows") or [])
            _add_chart(slide, spec["chart"], Inches(8.05), top, Inches(4.8), Inches(4.4))
        elif table.get("rows") is not None:
            _table(slide, Inches(0.34), top, Inches(12.6), Inches(6.55) - top,
                   table.get("headers") or ["No.", "Citation"], table.get("rows") or [])
        return

    if layout == "cards":
        _cards(slide, spec.get("cards") or [], Inches(0.34), top, Inches(12.6), Inches(6.55) - top)
        return

    if spec.get("chart"):
        _add_chart(slide, spec["chart"], Inches(0.34), top, Inches(12.6), Inches(2.2))
        top += Inches(2.2)
    if spec.get("bullets"):
        _bullets(slide, Inches(0.5), top, Inches(12.0), Inches(2.2),
                 spec["bullets"][:5], color=ink if not dark else WHITE, size=16)
    if spec.get("callout"):
        _callout(slide, Inches(0.34), Inches(6.25), Inches(12.4), Inches(0.5), spec["callout"], dark)


def _dark_chips(slide, cards: list[dict], l, t, w, h) -> None:
    n = max(len(cards), 1)
    cw = w / n
    for i, card in enumerate(cards):
        x = l + cw * i + Inches(0.06)
        _rail_box(slide, x, t, cw - Inches(0.12), h, _accent(i, card), DARK_CARD)
        _textbox(slide, x + Inches(0.18), t + Inches(0.12), cw - Inches(0.36), Inches(0.22),
                 card.get("meta") or "", size=10, color=BLUE)
        _textbox(slide, x + Inches(0.18), t + Inches(0.36), cw - Inches(0.36), Inches(0.4),
                 card.get("title") or "", size=14, bold=True, color=WHITE, font=DISPLAY)
        _textbox(slide, x + Inches(0.18), t + Inches(0.8), cw - Inches(0.36), h - Inches(0.92),
                 card.get("body") or "", size=11, color=RGBColor(0xD9, 0xD9, 0xD9))


def _dark_quad(slide, cards: list[dict], l, t, w, h) -> None:
    n = len(cards)
    cols = 2 if n != 3 else 3
    rows = 2 if n > 2 else 1
    cw = w / cols
    ch = h / max(rows, 1)
    for i, card in enumerate(cards[: cols * rows]):
        r, c = divmod(i, cols)
        x = l + cw * c + Inches(0.06)
        y = t + ch * r + Inches(0.06)
        _rail_box(slide, x, y, cw - Inches(0.12), ch - Inches(0.12), _accent(i, card), DARK_CARD)
        _textbox(slide, x + Inches(0.18), y + Inches(0.12), cw - Inches(0.36), Inches(0.38),
                 card.get("title") or "", size=16, bold=True, color=WHITE, font=DISPLAY)
        _textbox(slide, x + Inches(0.18), y + Inches(0.52), cw - Inches(0.36), ch - Inches(0.72),
                 card.get("body") or "", size=11, color=RGBColor(0xD9, 0xD9, 0xD9))


def _cards(slide, cards: list[dict], l, t, w, h) -> None:
    if not cards:
        return
    n = len(cards)
    cw = w / max(n, 1)
    for i, card in enumerate(cards):
        x = l + cw * i + Inches(0.06)
        ch = h
        _rail_box(slide, x, t, cw - Inches(0.12), ch, _accent(i, card), PAPER)
        _textbox(slide, x + Inches(0.18), t + Inches(0.12), cw - Inches(0.36), Inches(0.4),
                 card.get("title") or "", size=15, bold=True, color=INK, font=DISPLAY)
        _textbox(slide, x + Inches(0.18), t + Inches(0.55), cw - Inches(0.36), ch - Inches(1.0),
                 card.get("body") or "", size=12, color=BODY)
        if card.get("meta"):
            _textbox(slide, x + Inches(0.18), t + ch - Inches(0.38), cw - Inches(0.36), Inches(0.28),
                     card.get("meta") or "", size=10, color=ORANGE)


def _callout(slide, l, t, w, h, callout: dict, dark: bool) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.adjustments[0] = 0.08
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = NAVY if dark else RGBColor(0xF3, 0xE3, 0xC8)
    label = callout.get("label") or ""
    text = callout.get("text") or ""
    _textbox(slide, l + Inches(0.15), t + Inches(0.08), w - Inches(0.3), h - Inches(0.1),
             f"{label}. {text}".strip(), size=12, color=CREAM if dark else INK)


def _add_chart(slide, spec: dict, l, t, w, h) -> None:
    kind = spec.get("kind")
    data = spec.get("data") or []
    title = spec.get("title") or ""
    if kind not in {"people", "compare", "spine", "forest", "box", "flow", "house"} and title:
        _textbox(slide, l, t - Inches(0.28), w, Inches(0.28), title, size=11, color=MUTED)

    if kind == "forest":
        _forest(slide, data, l, t, w, h)
        return
    if kind == "box":
        _box(slide, data, l, t, w, h)
        return
    if kind == "people":
        _people(slide, data, l, t, w, h)
        return
    if kind == "compare":
        _compare(slide, data, l, t, w, h)
        return
    if kind == "spine":
        _spine(slide, data, l, t, w, h)
        return
    if kind == "flow":
        _flow(slide, data, l, t, w, h)
        return
    if kind == "house":
        _house(slide, spec, l, t, w, h)
        return
    if not data:
        return

    if kind == "pie":
        cd = CategoryChartData()
        cd.categories = [str(r.get("name", "")) for r in data]
        cd.add_series("Share", [_num(r.get("value")) for r in data])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, l, t, w, h, cd).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        return

    if kind == "line":
        series = spec.get("series") or ["value"]
        cd = CategoryChartData()
        cd.categories = [str(r.get("name", "")) for r in data]
        for s in series:
            cd.add_series(s, [_num(r.get(s)) for r in data])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, l, t, w, h, cd).chart
        chart.has_legend = True
        return

    if kind == "scatter":
        legend_w = min(int(w * 0.34), int(Inches(3.4)))
        chart_w = w - legend_w - Inches(0.12)
        cd = XyChartData()
        series = cd.add_series("Moves")
        for r in data:
            series.add_data_point(_num(r.get("x")), _num(r.get("y")))
        slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, l, t, chart_w, h, cd)
        lines = [
            f"{r.get('name')}: {spec.get('xLabel') or 'Feasibility'} {r.get('x')}, {spec.get('yLabel') or 'Impact'} {r.get('y')}."
            for r in data
        ]
        _bullets(slide, l + chart_w + Inches(0.12), t, legend_w, h, lines, size=12)
        return

    if kind == "diverging":
        cd = CategoryChartData()
        cd.categories = [str(r.get("name", "")) for r in data]
        cd.add_series("Matches evidence", [max(_num(r.get("value")), 0) for r in data])
        cd.add_series("Perception gap", [min(_num(r.get("value")), 0) for r in data])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, cd).chart
        chart.has_legend = True
        return

    # bar / funnel default
    cd = CategoryChartData()
    cd.categories = [str(r.get("name", "")) for r in data]
    cd.add_series(spec.get("unit") or "Value", [_num(r.get("value")) for r in data])
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)


def _num(value: Any) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return 0.0


def _forest(slide, rows: list[dict], l, t, w, h) -> None:
    """Editable forest plot: lines + diamonds, not a flattened image."""
    if not rows:
        return
    nums = []
    for r in rows:
        nums.extend([_num(r.get("low")), _num(r.get("high")), _num(r.get("hr"))])
    lo, hi = min(0.55, min(nums) - 0.04), max(1.15, max(nums) + 0.04)
    plot_l = l + Inches(2.3)
    plot_w = w - Inches(3.9)
    row_h = int((h - Inches(0.28)) / max(len(rows), 1))

    def x_at(v: float) -> int:
        return int(plot_l + ((v - lo) / (hi - lo)) * plot_w)

    # Null line at 1.0
    null_x = x_at(1.0)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, null_x, t, Emu(12700), h - Inches(0.2))
    line.fill.solid()
    line.fill.fore_color.rgb = COPPER
    line.line.fill.background()
    _textbox(slide, null_x - Inches(0.55), t - Inches(0.22), Inches(1.2), Inches(0.22),
             "null 1.0", size=9, color=COPPER, align=PP_ALIGN.CENTER)
    _textbox(slide, plot_l, t + h - Inches(0.22), Inches(2.4), Inches(0.22),
             "Favours intervention", size=9, color=TEAL)
    _textbox(slide, plot_l + plot_w - Inches(2.4), t + h - Inches(0.22), Inches(2.4), Inches(0.22),
             "Favours comparator", size=9, color=CRIMSON, align=PP_ALIGN.RIGHT)

    for i, r in enumerate(rows):
        y = t + row_h * i + int(row_h * 0.35)
        _textbox(slide, l, y - Inches(0.12), Inches(2.2), Inches(0.45),
                 f"{r.get('name', '')}\n{r.get('stream', '')} · {r.get('grade', '')}",
                 size=10, color=INK)
        x1, x2, xm = x_at(_num(r.get("low"))), x_at(_num(r.get("high"))), x_at(_num(r.get("hr")))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, min(x1, x2), y, max(abs(x2 - x1), Emu(12700)), Emu(19050))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, xm - Emu(69850), y - Emu(50800), Emu(139700), Emu(139700))
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = COPPER
        diamond.line.fill.background()
        hr, low, high = _num(r.get("hr")), _num(r.get("low")), _num(r.get("high"))
        _textbox(
            slide,
            l + w - Inches(1.55),
            y - Inches(0.14),
            Inches(1.5),
            Inches(0.4),
            f"{hr:.2f} ({low:.2f}–{high:.2f})",
            size=9,
            color=MUTED,
            align=PP_ALIGN.RIGHT,
        )


def _box(slide, rows: list[dict], l, t, w, h) -> None:
    if not rows:
        return
    col_w = int(w / len(rows))
    axis_top, axis_bot = t + Inches(0.15), t + h - Inches(0.45)

    def y_at(v: float) -> int:
        return int(axis_bot - (v / 10.0) * (axis_bot - axis_top))

    for i, r in enumerate(rows):
        cx = l + col_w * i + col_w // 2
        y_min, y_q1 = y_at(_num(r.get("min"))), y_at(_num(r.get("q1")))
        y_med, y_q3 = y_at(_num(r.get("median"))), y_at(_num(r.get("q3")))
        y_max = y_at(_num(r.get("max")))
        whisker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, min(y_max, y_min), Emu(12700), abs(y_min - y_max))
        whisker.fill.solid()
        whisker.fill.fore_color.rgb = NAVY
        whisker.line.fill.background()
        box_top, box_bot = min(y_q3, y_q1), max(y_q3, y_q1)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - Inches(0.22), box_top, Inches(0.44), max(box_bot - box_top, Emu(25400)))
        box.fill.solid()
        box.fill.fore_color.rgb = TEAL
        box.line.color.rgb = NAVY
        med = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - Inches(0.22), y_med, Inches(0.44), Emu(19050))
        med.fill.solid()
        med.fill.fore_color.rgb = CREAM
        med.line.fill.background()
        _textbox(slide, l + col_w * i, t + h - Inches(0.4), col_w, Inches(0.4),
                 str(r.get("name", "")), size=10, color=INK, align=PP_ALIGN.CENTER)


def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _people(slide, rows: list[dict], l, t, w, h) -> None:
    """100-dot event grids: comparator vs intervention, plus NNT card."""
    if not rows:
        return
    row = rows[0]
    control = int(round(_num(row.get("control"))))
    treat = int(round(_num(row.get("treat"))))
    arr_raw = row.get("arr")
    arr = int(round(_num(arr_raw) if arr_raw not in (None, "") else max(0, control - treat)))
    panel_w = int(w * 0.34)
    nnt_w = w - panel_w * 2 - Inches(0.3)
    _people_grid(slide, l, t, panel_w, h, control, 0, str(row.get("control_label") or "Comparator"), CRIMSON)
    _people_grid(slide, l + panel_w + Inches(0.12), t, panel_w, h, treat, arr, str(row.get("treat_label") or "Intervention"), TEAL)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l + panel_w * 2 + Inches(0.24), t, nnt_w, h - Inches(0.1))
    card.adjustments[0] = 0.08
    _solid(card, NAVY)
    _textbox(slide, l + panel_w * 2 + Inches(0.4), t + Inches(0.15), nnt_w - Inches(0.3), Inches(0.28),
             "NNT", size=11, color=ORANGE)
    _textbox(slide, l + panel_w * 2 + Inches(0.4), t + Inches(0.45), nnt_w - Inches(0.3), Inches(1.1),
             str(row.get("nnt") or "—"), size=44, bold=True, color=WHITE, font=DISPLAY)
    _textbox(slide, l + panel_w * 2 + Inches(0.4), t + Inches(1.6), nnt_w - Inches(0.3), Inches(1.6),
             f"Treat {row.get('nnt') or '—'} to prevent 1 event"
             + (f" over {row.get('horizon')}." if row.get("horizon") else ".")
             + f" PMID {row.get('pmid') or '—'}.",
             size=13, color=WHITE)


def _people_grid(slide, l, t, w, h, events: int, saved: int, label: str, accent: RGBColor) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h - Inches(0.1))
    panel.adjustments[0] = 0.06
    panel.fill.solid()
    panel.fill.fore_color.rgb = PAPER
    panel.line.color.rgb = RGBColor(0xE0, 0xD8, 0xC8)
    _textbox(slide, l + Inches(0.12), t + Inches(0.08), w - Inches(0.2), Inches(0.28),
             label, size=11, color=MUTED)
    _textbox(slide, l + Inches(0.12), t + Inches(0.32), w - Inches(0.2), Inches(0.4),
             f"{events} / 100", size=22, bold=True, color=INK)
    origin_l = l + Inches(0.16)
    origin_t = t + Inches(0.8)
    cell = min(int((w - Inches(0.32)) / 10), int((h - Inches(1.15)) / 10))
    events = max(0, min(100, events))
    saved = max(0, min(100 - events, saved))
    for i in range(100):
        r, c = divmod(i, 10)
        x = origin_l + c * cell
        y = origin_t + r * cell
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + 2, y + 2, cell - 4, cell - 4)
        if i < events:
            color = CRIMSON
        elif i < events + saved:
            color = TEAL
        else:
            color = RGBColor(0xD4, 0xCE, 0xC3)
        _solid(dot, color)
    _textbox(slide, l + Inches(0.12), t + h - Inches(0.38), w - Inches(0.2), Inches(0.24),
             "event" if accent == CRIMSON else "event + avoided", size=9, color=MUTED)


def _compare(slide, rows: list[dict], l, t, w, h) -> None:
    if not rows:
        return
    row = rows[0]
    left, right = _num(row.get("left")), _num(row.get("right"))
    mx = max(left, right, 1)
    col_w = int(w * 0.28)
    mid_w = w - col_w * 2 - Inches(0.3)
    bar_h = h - Inches(1.3)

    def panel(x, value, label, color):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, t, col_w, h - Inches(0.1))
        box.adjustments[0] = 0.06
        box.fill.solid()
        box.fill.fore_color.rgb = PAPER
        box.line.color.rgb = RGBColor(0xE0, 0xD8, 0xC8)
        _textbox(slide, x + Inches(0.1), t + Inches(0.08), col_w - Inches(0.2), Inches(0.28),
                 str(label), size=11, color=MUTED)
        _textbox(slide, x + Inches(0.1), t + Inches(0.36), col_w - Inches(0.2), Inches(0.45),
                 str(int(value) if value == int(value) else value), size=28, bold=True, color=INK)
        fill_h = max(int(bar_h * (value / mx)), Emu(25400))
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + int(col_w * 0.35),
            t + Inches(0.9) + (bar_h - fill_h),
            int(col_w * 0.28),
            fill_h,
        )
        _solid(bar, color)

    panel(l, left, row.get("left_label") or "Comparator", CRIMSON)
    mid = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l + col_w + Inches(0.15), t, mid_w, h - Inches(0.1))
    mid.adjustments[0] = 0.06
    _solid(mid, RGBColor(0xF3, 0xE3, 0xC8))
    delta = row.get("delta") if row.get("delta") not in (None, "") else abs(left - right)
    _textbox(slide, l + col_w + Inches(0.3), t + Inches(0.15), mid_w - Inches(0.3), Inches(0.28),
             "Difference", size=11, color=COPPER)
    _textbox(slide, l + col_w + Inches(0.3), t + Inches(0.45), mid_w - Inches(0.3), Inches(0.6),
             str(delta), size=32, bold=True, color=INK)
    _textbox(slide, l + col_w + Inches(0.3), t + Inches(1.15), mid_w - Inches(0.3), Inches(2.2),
             f"{row.get('claim') or ''}  PMID {row.get('pmid') or '—'}.", size=12, color=MUTED)
    panel(l + col_w + mid_w + Inches(0.3), right, row.get("right_label") or "Intervention", TEAL)


def _spine(slide, rows: list[dict], l, t, w, h) -> None:
    if not rows:
        return
    labels = [("science", "1. Science"), ("means", "2. Means"), ("barrier", "3. Barrier"),
              ("execute", "4. Execution"), ("measure", "5. We measure")]
    col_w = int(w / 5)
    _textbox(slide, l, t, w, Inches(0.22), "Science → means → barrier → execution → we measure", size=10, color=COPPER)
    row_h = int((h - Inches(0.28)) / max(min(len(rows), 3), 1))
    for i, row in enumerate(rows[:3]):
        y = t + Inches(0.28) + row_h * i
        for j, (key, label) in enumerate(labels):
            x = l + col_w * j + Inches(0.04)
            cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w - Inches(0.08), row_h - Inches(0.08))
            cell.adjustments[0] = 0.08
            if key == "execute":
                _solid(cell, NAVY)
                ink, muted = CREAM, COPPER
            elif key == "measure":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xEE)
                cell.line.fill.background()
                ink, muted = INK, TEAL
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER
                cell.line.color.rgb = RGBColor(0xE0, 0xD8, 0xC8)
                ink, muted = INK, MUTED
            heading = label if key != "science" else f"{label} · {row.get('name', '')}"
            _textbox(slide, x + Inches(0.08), y + Inches(0.06), col_w - Inches(0.2), Inches(0.28),
                     heading, size=9, color=muted)
            text = str(row.get(key) or "—")
            if key == "science" and row.get("pmid"):
                text = f"{text}  PMID {row.get('pmid')}."
            _textbox(slide, x + Inches(0.08), y + Inches(0.32), col_w - Inches(0.2), row_h - Inches(0.42),
                     text, size=10, color=ink)


def _flow(slide, rows: list[dict], l, t, w, h) -> None:
    if not rows:
        return
    n = len(rows)
    cw = w / n
    for i, row in enumerate(rows):
        x = l + cw * i + Inches(0.08)
        box_w = cw - Inches(0.28)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, t, box_w, h - Inches(0.08))
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = PAPER
        box.line.color.rgb = RGBColor(0xE0, 0xD8, 0xC8)
        _textbox(slide, x + Inches(0.12), t + Inches(0.1), box_w - Inches(0.2), Inches(0.28),
                 str(i + 1), size=11, color=COPPER)
        _textbox(slide, x + Inches(0.12), t + Inches(0.38), box_w - Inches(0.2), Inches(0.55),
                 str(row.get("name") or ""), size=13, bold=True, color=INK)
        _textbox(slide, x + Inches(0.12), t + Inches(0.95), box_w - Inches(0.2), h - Inches(1.15),
                 str(row.get("detail") or row.get("claim") or ""), size=12, color=MUTED)
        if i < n - 1:
            _textbox(slide, x + box_w - Inches(0.05), t + (h / 2) - Inches(0.18), Inches(0.28), Inches(0.36),
                     "→", size=16, bold=True, color=COPPER, align=PP_ALIGN.CENTER)


def _house(slide, spec: dict, l, t, w, h) -> None:
    rows = spec.get("data") or []
    _textbox(slide, l, t, w, Inches(0.22), "THEME", size=10, bold=True, color=ORANGE)
    _textbox(slide, l, t + Inches(0.22), w, Inches(0.45),
             spec.get("title") or "", size=16, bold=True, color=INK, font=DISPLAY)
    if not rows:
        return
    n = len(rows)
    cw = w / n
    y = t + Inches(0.78)
    ch = h - Inches(0.85)
    for i, row in enumerate(rows):
        x = l + cw * i + Inches(0.06)
        _rail_box(slide, x, y, cw - Inches(0.12), ch, _accent(i), PAPER)
        _textbox(slide, x + Inches(0.16), y + Inches(0.1), cw - Inches(0.34), Inches(0.22),
                 str(row.get("ref") or "Pillar"), size=10, bold=True, color=ORANGE)
        _textbox(slide, x + Inches(0.16), y + Inches(0.34), cw - Inches(0.34), Inches(0.4),
                 str(row.get("name") or ""), size=15, bold=True, color=INK, font=DISPLAY)
        _textbox(slide, x + Inches(0.16), y + Inches(0.78), cw - Inches(0.34), ch - Inches(1.2),
                 str(row.get("line") or ""), size=12, color=BODY)
        _textbox(slide, x + Inches(0.16), y + ch - Inches(0.38), cw - Inches(0.34), Inches(0.32),
                 str(row.get("proof") or ""), size=11, color=MUTED)


def _appendix_bibliography(prs, blank, pack: dict) -> None:
    records = (pack.get("evidence") or {}).get("records") or pack.get("dashboard", {}).get("citations") or []
    if not records:
        return
    slide = prs.slides.add_slide(blank)
    _fill_slide(slide, CREAM)
    _textbox(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.32),
             "APPENDIX  ·  VALIDATED SOURCES", size=11, color=COPPER)
    lead = (pack.get("evidence") or {}).get("lead") or {}
    _textbox(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.55),
             "Bibliography — campaign lead is the first row", size=22, bold=True, color=INK)
    if lead.get("statement"):
        _textbox(slide, Inches(0.55), Inches(1.1), Inches(12.2), Inches(0.7),
                 lead["statement"], size=13, color=MUTED)
    headers = ["Short", "Citation", "PMID", "DOI", "Grade", "What we may claim"]
    rows = [
        [
            r.get("short") or "",
            r.get("citation") or "",
            str(r.get("pmid") or ""),
            str(r.get("doi") or ""),
            r.get("grade") or "",
            r.get("claim_permitted") or "",
        ]
        for r in records
    ]
    _table(slide, Inches(0.35), Inches(1.85), Inches(12.6), Inches(5.3), headers, rows)


def _appendix_interventions(prs, blank, pack: dict) -> None:
    items = pack.get("interventions") or []
    if not items:
        return
    slide = prs.slides.add_slide(blank)
    _fill_slide(slide, CREAM)
    _textbox(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.32),
             "APPENDIX  ·  INTERVENTION BOARD", size=11, color=COPPER)
    _textbox(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.5),
             "Five moves — owners can edit effort, copy, and kill-criteria", size=22, bold=True, color=INK)
    headers = ["Move", "Promise", "COM-B lever", "Segment", "Impact", "Feasibility", "Kill-criterion"]
    rows = [
        [i.get("name", ""), i.get("promise", ""), i.get("lever", ""), i.get("segment", ""),
         str(i.get("impact", "")), str(i.get("feasibility", "")), i.get("kill", "")]
        for i in items
    ]
    _table(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.7), headers, rows)


def _appendix_dashboard(prs, blank, pack: dict) -> None:
    dash = pack.get("dashboard") or {}
    kpis = dash.get("kpis") or []
    if not kpis:
        return
    slide = prs.slides.add_slide(blank)
    _fill_slide(slide, CREAM)
    _textbox(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.32),
             "APPENDIX  ·  MEASUREMENT ROOM", size=11, color=COPPER)
    _textbox(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.5),
             "Starting KPIs — replace illustrative numbers before client lock", size=22, bold=True, color=INK)
    headers = ["KPI", "Now", "Target", "Unit", "Type"]
    rows = [[k.get("label", ""), str(k.get("value", "")), str(k.get("target", "")),
             k.get("unit", ""), k.get("tone", "")] for k in kpis]
    _table(slide, Inches(0.55), Inches(1.25), Inches(6.2), Inches(3.2), headers, rows)

    revenue = dash.get("revenue") or []
    if revenue:
        cd = CategoryChartData()
        cd.categories = [str(r.get("name", "")) for r in revenue]
        for s in ("revenue", "initiation", "conviction"):
            if s in revenue[0]:
                cd.add_series(s, [_num(r.get(s)) for r in revenue])
        slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(7.0), Inches(1.25), Inches(5.7), Inches(3.4), cd)

    gov = dash.get("governance") or []
    if gov:
        _table(slide, Inches(0.55), Inches(4.7), Inches(12.2), Inches(2.3),
               ["Cadence", "Forum", "Looks at"],
               [[g.get("cadence", ""), g.get("forum", ""), g.get("looksAt", "")] for g in gov])


def _appendix_edit_guide(prs, blank, meta: dict, doctrine: dict) -> None:
    slide = prs.slides.add_slide(blank)
    _fill_slide(slide, INK)
    _textbox(slide, Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.32),
             "HOW TO EDIT THIS DECK", size=11, color=COPPER)
    _textbox(slide, Inches(0.55), Inches(0.8), Inches(12.2), Inches(0.7),
             "Every object is a native PowerPoint shape, chart, or table.", size=26, bold=True, color=CREAM)
    _bullets(slide, Inches(0.7), Inches(1.8), Inches(12.0), Inches(4.8), [
        "Click any title, bullet, or table cell and type. Fonts are Calibri so they travel.",
        "Bar, line, pie, and scatter charts are Office charts — right-click → Edit Data to change numbers.",
        "Forest and box plots are grouped shapes. Drag a diamond or box to restyle; replace values in the labels.",
        "Speaker notes under each slide carry the narrative and any evidence caveats.",
        "MLR: no promotional use until medical-legal-regulatory clears claims and field scripts.",
        f"Brand: {meta.get('brand', '—')}  ·  Doctrine: {doctrine.get('name', meta.get('doctrine', '—'))}  ·  Generated: {meta.get('generatedAt', '—')}",
    ], color=CREAM, size=16)


