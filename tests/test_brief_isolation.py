"""One brief in, one brief out — never mix a template or a leftover pack."""

import io
import json

from pptx import Presentation
from pptx.util import Inches, Pt

from director_api.app import app
from director_api.evidence import resolve_evidence
from director_api.extract import ExtractedBrief, extract_one, is_design_template, merge_into_brief, partition_uploads
from director_api.generate import generate_pack
from fastapi.testclient import TestClient

client = TestClient(app)

HELIX = b"""brand: HelixOne
therapy_area: Endocrinology
indication: Type 2 diabetes
market: India
product: HelixOne
hcp_insights:
  - Endocrinologists accept the science but start late
business_goal: Move initiation to the first eligible visit
"""


def _fake_omnicom_template() -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1))
    box.text_frame.text = "MEDICAL STRATEGY DECK · SEMANEXT®"
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(10), Inches(0.4))
    box2.text_frame.text = "© 2026 Omnicom Group Inc. All rights reserved."
    box3 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(10), Inches(0.8))
    box3.text_frame.paragraphs[0].font.size = Pt(14)
    box3.text_frame.text = "LONG-TERM GLP-1 THERAPY and Ozempic competitive claims"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_omnicom_pptx_is_a_design_template_not_a_brief():
    payload = _fake_omnicom_template()
    name = "Semanext_LongTerm_GLP1_Strategy_Omnicom_Template_v0.2.pptx"
    extracted = extract_one(name, payload)
    assert is_design_template(name, extracted.text)
    templates, briefs = partition_uploads(
        [
            (name, payload, "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            ("helix.yaml", HELIX, "text/yaml"),
        ]
    )
    assert templates and briefs
    brief = merge_into_brief([extract_one("helix.yaml", HELIX, "text/yaml")])
    assert brief.brand == "HelixOne"
    assert "Semanext" not in (brief.brand + brief.therapy_area)


def test_generate_does_not_mix_template_or_cardioshield_into_helixone():
    payload = _fake_omnicom_template()
    name = "Omnicom_Template.pptx"
    files = [
        ("files", (name, payload, "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
        ("files", ("helix.yaml", HELIX, "text/yaml")),
    ]
    stale = {
        "brand": "CardioShield",
        "therapy_area": "Cardiology - chronic heart failure",
        "indication": "HFrEF",
        "source_files": ["old-cardio.yaml"],
    }
    res = client.post(
        "/api/generate",
        files=files,
        data={"brief_json": json.dumps(stale), "pasted": ""},
    )
    assert res.status_code == 200
    pack = res.json()
    assert pack["meta"]["brand"] == "HelixOne"
    assert pack["meta"]["therapyArea"] == "Endocrinology"
    blob = " ".join(
        [
            pack["meta"]["brand"],
            pack["slides"][0]["title"],
            pack["slides"][0].get("subtitle") or "",
            *(s.get("title") or "" for s in pack["slides"]),
            *(s.get("narrative") or "" for s in pack["slides"]),
            (pack.get("evidence") or {}).get("lead", {}).get("statement") or "",
        ]
    )
    assert "HelixOne" in blob
    assert "CardioShield" not in blob
    assert "Semanext" not in blob
    assert "PARADIGM-HF" not in blob
    assert "PIONEER-HF" not in blob
    ids = {r["id"] for r in pack["evidence"]["records"]}
    assert "paradigm-hf-2014" not in ids
    assert "pioneer-hf-2019" not in ids
    assert "keynote-189-2018" not in ids
    assert pack["meta"].get("templateFiles")


def test_endocrinology_brief_does_not_inherit_heart_failure_catalog():
    brief = ExtractedBrief(
        brand="HelixOne",
        therapy_area="Endocrinology",
        indication="Type 2 diabetes",
        market="India",
        business_goal="Move initiation to the first eligible visit",
        hcp_insights=["Endocrinologists accept the science but start late"],
        raw_text="early initiation in India",
    )
    ledger = resolve_evidence(brief, pubmed=False)
    ids = {r["id"] for r in ledger["records"]}
    assert "pioneer-hf-2019" not in ids
    assert "trivandrum-hf-2015" not in ids
    assert "paradigm-hf-2014" not in ids
    pack = generate_pack(brief, pubmed=False)
    assert pack["meta"]["brand"] == "HelixOne"
    assert pack["evidence"]["records"] == []
    blob = json.dumps(pack).lower()
    assert "keynote" not in blob
    assert "pembrolizumab" not in blob
    assert "sacubitril" not in blob
