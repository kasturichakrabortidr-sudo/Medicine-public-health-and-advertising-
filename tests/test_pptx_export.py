import io
import zipfile

from pptx import Presentation

from director_api.app import _brief_from_mapping, app
from director_api.generate import generate_pack
from director_api.pptx_export import filename_for, pack_to_pptx
from fastapi.testclient import TestClient
from medicomarketing_agent.config import load_brief

client = TestClient(app)


def _demo_pack():
    return generate_pack(_brief_from_mapping(load_brief("examples/brief.example.yaml")), mode="demo")


def test_pptx_is_editable_office_file():
    pack = _demo_pack()
    data = pack_to_pptx(pack)
    assert data[:2] == b"PK"
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        assert "ppt/presentation.xml" in names
        assert any(n.startswith("ppt/slides/slide") for n in names)
        assert any(n.startswith("ppt/charts/") for n in names)

    prs = Presentation(io.BytesIO(data))
    # Strategy slides + bibliography + interventions + dashboard + edit guide
    assert len(prs.slides) == len(pack["slides"]) + 4
    title = prs.slides[0]
    texts = [shape.text_frame.text for shape in title.shapes if shape.has_text_frame]
    assert any("CardioShield" in t for t in texts)
    assert filename_for(pack) == "CardioShield-strategy-deck.pptx"
    meaning_idx = next(i for i, s in enumerate(pack["slides"]) if s["id"] == "science-meaning")
    meaning_texts = [shape.text_frame.text for shape in prs.slides[meaning_idx].shapes if shape.has_text_frame]
    joined = " ".join(meaning_texts)
    assert "21" in joined
    assert "25176015" in joined
    execute_idx = next(i for i, s in enumerate(pack["slides"]) if s["id"] == "science-execute")
    execute_texts = [shape.text_frame.text for shape in prs.slides[execute_idx].shapes if shape.has_text_frame]
    execute_joined = " ".join(execute_texts)
    assert "First-Touch" in execute_joined or "first-eligible" in execute_joined.lower()


def test_export_endpoints():
    res = client.get("/api/export/pptx")
    assert res.status_code == 200
    assert "presentationml" in res.headers["content-type"]
    assert res.headers["content-disposition"].endswith('.pptx"')
    assert res.content[:2] == b"PK"
    Presentation(io.BytesIO(res.content))

    pack = client.get("/api/demo").json()
    posted = client.post("/api/export/pptx", json=pack)
    assert posted.status_code == 200
    assert posted.content[:2] == b"PK"


def test_export_rejects_empty_pack():
    res = client.post("/api/export/pptx", json={"meta": {}})
    assert res.status_code == 400
