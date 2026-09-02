from pathlib import Path

from director_api.extract import extract_one, merge_into_brief


def test_yaml_brief_fields():
    payload = b"""
brand: CardioShield
therapy_area: Cardiology - chronic heart failure
hcp_insights:
  - Field notes - cost is the veto
"""
    extracted = extract_one("brief.yaml", payload, "text/yaml")
    brief = merge_into_brief([extracted])
    assert brief.brand == "CardioShield"
    assert "Cardiology" in brief.therapy_area
    assert brief.hcp_insights[0].startswith("Field notes")


def test_json_brief():
    extracted = extract_one(
        "brief.json",
        b'{"brand":"Lumen","therapy_area":"Oncology"}',
        "application/json",
    )
    brief = merge_into_brief([extracted])
    assert brief.brand == "Lumen"
    assert brief.therapy_area == "Oncology"


def test_csv_and_markdown_heading():
    csv_file = extract_one("notes.csv", b"theme,detail\ncost,OOP burden\n", "text/csv")
    md = extract_one(
        "notes.md",
        b"# Brand\nHelix\n## Therapy Area\nNeurology\n",
        "text/markdown",
    )
    brief = merge_into_brief([csv_file, md])
    assert "Helix" in brief.brand
    assert "Neurology" in brief.therapy_area
    assert "cost" in brief.raw_text


def test_html_and_rtf_strip_markup():
    html = extract_one("a.html", b"<html><body><p>Brand: Nimbus</p></body></html>", "text/html")
    rtf = extract_one("a.rtf", b"{\\rtf1 Therapy area: Dermatology}", "application/rtf")
    assert "Nimbus" in html.text
    assert "<p>" not in html.text
    assert "Dermatology" in rtf.text


def test_image_is_noted_not_ocrd():
    extracted = extract_one("scan.png", b"\x89PNG\r\n", "image/png")
    assert extracted.notes
    assert "scan.png" in extracted.text


def test_docx_and_pptx_extract_text():
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches
    import io

    doc = Document()
    doc.add_heading("Brand", level=1)
    doc.add_paragraph("Nimbus")
    buf = io.BytesIO()
    doc.save(buf)
    brief = merge_into_brief([extract_one("b.docx", buf.getvalue())])
    assert "Nimbus" in brief.raw_text

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "Therapy area: Immunology"
    pbuf = io.BytesIO()
    prs.save(pbuf)
    ppt = extract_one("deck.pptx", pbuf.getvalue())
    assert "Immunology" in ppt.text


def test_labeled_prose_fills_hcp_habit():
    brief = merge_into_brief(
        [],
        pasted=(
            "Brand: Helix\nTherapy area: Oncology\nIndication: metastatic NSCLC\n"
            "HCP habit: Most wait until the patient is stable in clinic\n"
            "Business goal: Move first-eligible starts this quarter\n"
            "MLR: No comparative claim without a numbered paper\n"
        ),
    )
    assert brief.brand == "Helix"
    assert "Oncology" in brief.therapy_area
    assert brief.hcp_insights
    assert "stable" in brief.hcp_insights[0].lower()
    assert "first-eligible" in brief.business_goal.lower()
    assert any("comparative" in c.lower() for c in brief.constraints)


def test_alex_creative_brief_is_read():
    raw = (Path(__file__).resolve().parent / "fixtures" / "alex_creative_brief.txt").read_text(
        encoding="utf-8"
    )
    brief = merge_into_brief([], pasted=raw)
    assert brief.brand == "Alex"
    assert "cough" in (brief.indication or brief.therapy_area).lower()
    assert "loss of rank" not in (brief.therapy_area or "").lower()
    assert any("similar" in line.lower() for line in brief.hcp_insights)
    assert any("Zedex" in c or "Grilinctus" in c for c in brief.competitors)
    assert "Alex" in (brief.business_goal or "")


def test_novartis_portfolio_brief_leads_with_sybrava():
    raw = (Path(__file__).resolve().parent / "fixtures" / "novartis_cv_brief.txt").read_text(
        encoding="utf-8"
    )
    brief = merge_into_brief([], pasted=raw)
    assert brief.brand == "Sybrava"
    assert "Inclisiran" in (brief.product or "")
    assert "Cardio" in (brief.therapy_area or "") or "LDL" in (brief.indication or "")
    assert "Pelacarsen" not in " ".join(brief.hcp_insights)


def test_table_style_creative_brief_is_read():
    brief = merge_into_brief(
        [],
        pasted=(
            "What objective are we aiming to achieve? | Grow unaided recall of Nimbus\n"
            "What is their Current belief / behavior? | Doctors treat every brand as interchangeable\n"
            "Who’s the target audience (types of HCPs)? | Dermatologists\n"
        ),
    )
    assert brief.brand == "Nimbus" or "Nimbus" in (brief.business_goal or "")
    assert any("interchangeable" in line.lower() for line in brief.hcp_insights)
    assert any("Dermatolog" in s for s in brief.target_specialties)


def test_alex_filename_fills_brand_when_prose_is_thin():
    from director_api.extract import ExtractedFile

    extracted = ExtractedFile(
        filename="Glenmark-Alex_creative brief format_mccann health.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        suffix=".docx",
        text="Creative brief. Dry cough category. Rival visibility from Zedex.",
    )
    brief = merge_into_brief([extracted])
    assert brief.brand == "Alex"
    assert any("Zedex" in c for c in brief.competitors)


TRIVORA = (Path(__file__).resolve().parent / "fixtures" / "trivora_client_brief.txt").read_text(
    encoding="utf-8"
)


def test_trivora_client_brief_is_not_swallowed_into_brand():
    brief = merge_into_brief([], pasted=TRIVORA)
    assert brief.brand == "Trivora-NB Smartules"
    assert "Formoterol" in brief.product or "formoterol" in brief.product
    assert "Glycopyrronium" in brief.product
    assert "COPD" in brief.therapy_area
    assert "first-line" in (brief.business_goal or "").lower()
    assert "20%" in (brief.business_goal or "")
    assert any("free-mix" in c.lower() for c in brief.competitors)
    assert any("GOLD" in g for g in brief.guidelines)
    assert len(brief.hcp_insights) >= 5
    blob = " ".join(brief.hcp_insights).lower()
    assert "rescue or step-up" in blob or "when dual" in blob
    assert "mix it myself" in blob
    assert "consultant physicians, metro and tier-1)" not in blob
    assert "Molecule" not in brief.brand
    assert len(brief.brand) < 60
    assert "GOLD 2026" not in (brief.indication or "")
    assert "STATED POPULATION" not in (brief.indication or "")
